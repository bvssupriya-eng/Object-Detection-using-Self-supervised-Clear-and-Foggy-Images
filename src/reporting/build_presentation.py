"""Generate a professional project presentation as a PowerPoint file."""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


PROJECT_ROOT = Path(__file__).resolve().parents[2]
REPORT_ROOT = PROJECT_ROOT / "report"
RESULTS_ROOT = PROJECT_ROOT / "results"
OUTPUT_PATH = REPORT_ROOT / "final_project_presentation.pptx"


BG = RGBColor(8, 17, 29)
PANEL = RGBColor(17, 28, 43)
PANEL_ALT = RGBColor(11, 41, 66)
PANEL_SOFT = RGBColor(23, 36, 54)
ACCENT = RGBColor(93, 224, 230)
ACCENT_2 = RGBColor(247, 181, 56)
TEXT = RGBColor(245, 247, 250)
MUTED = RGBColor(191, 203, 214)
SUCCESS = RGBColor(97, 214, 135)


def set_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.0), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.size = Pt(26)
    run.font.bold = True
    run.font.color.rgb = TEXT
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(12.0), Inches(0.4))
        sub_tf = sub_box.text_frame
        sub_p = sub_tf.paragraphs[0]
        sub_run = sub_p.add_run()
        sub_run.text = subtitle
        sub_run.font.size = Pt(11)
        sub_run.font.color.rgb = ACCENT


def add_footer(slide, text: str = "MachineVision Project") -> None:
    box = slide.shapes.add_textbox(Inches(0.6), Inches(7.0), Inches(12.0), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    run.font.size = Pt(9)
    run.font.color.rgb = MUTED


def add_panel(slide, left, top, width, height, color=PANEL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(44, 61, 84)
    shape.line.width = Pt(1)
    return shape


def add_top_band(slide, label: str) -> None:
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.16),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = ACCENT
    band.line.fill.background()
    add_textbox(slide, label, Inches(10.7), Inches(0.18), Inches(2.0), Inches(0.2), size=9, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)


def add_bullets(slide, items: list[str], left, top, width, height, font_size=18) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT
        p.space_after = Pt(10)
        p.bullet = True


def add_textbox(slide, text: str, left, top, width, height, size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return box


def add_image(slide, path: Path, left, top, width=None, height=None):
    if path.exists():
        slide.shapes.add_picture(str(path), left, top, width=width, height=height)


def add_metric_box(slide, title: str, value: str, left, top, width, height, accent=ACCENT):
    shape = add_panel(slide, left, top, width, height, color=PANEL_ALT)
    shape.line.color.rgb = accent
    add_textbox(slide, title, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.22), size=11, color=MUTED, bold=True)
    add_textbox(slide, value, left + Inches(0.15), top + Inches(0.33), width - Inches(0.3), Inches(0.3), size=22, color=TEXT, bold=True)


def add_kpi_row(slide, items: list[tuple[str, str]], top: float) -> None:
    left = Inches(0.82)
    width = Inches(2.85)
    gap = Inches(0.23)
    for idx, (label, value) in enumerate(items):
        accent = [ACCENT, SUCCESS, ACCENT_2, RGBColor(255, 142, 142)][idx % 4]
        add_metric_box(slide, label, value, left, Inches(top), width, Inches(0.95), accent=accent)
        left += width + gap


def add_table(slide, data: pd.DataFrame, left, top, width, height) -> None:
    rows, cols = data.shape[0] + 1, data.shape[1]
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    for col_idx, column in enumerate(data.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(column)
        cell.fill.solid()
        cell.fill.fore_color.rgb = PANEL_ALT
        for paragraph in cell.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.bold = True
                run.font.color.rgb = TEXT
                run.font.size = Pt(12)
    for row_idx, row in enumerate(data.itertuples(index=False), start=1):
        for col_idx, value in enumerate(row):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(value)
            cell.fill.solid()
            cell.fill.fore_color.rgb = PANEL if row_idx % 2 == 1 else RGBColor(21, 34, 52)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.color.rgb = TEXT
                    run.font.size = Pt(11)


def add_workflow_box(slide, label: str, left, top, width=Inches(1.35), height=Inches(0.65), color=PANEL_ALT):
    shape = add_panel(slide, left, top, width, height, color=color)
    shape.adjustments[0] = 0.08
    add_textbox(slide, label, left + Inches(0.05), top + Inches(0.08), width - Inches(0.1), height - Inches(0.16), size=11, bold=True, align=PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2)
    line.line.end_arrowhead = True
    return line


def add_stage_box(slide, title: str, subtitle: str, left, top, width, height, color=PANEL_ALT):
    shape = add_panel(slide, left, top, width, height, color=color)
    add_textbox(slide, title, left + Inches(0.1), top + Inches(0.11), width - Inches(0.2), Inches(0.26), size=14, color=TEXT, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, subtitle, left + Inches(0.12), top + Inches(0.42), width - Inches(0.24), Inches(0.46), size=10, color=MUTED, align=PP_ALIGN.CENTER)
    return shape


def add_architecture(slide, model_name: str, desc: str, accent: RGBColor) -> None:
    add_title(slide, f"{model_name} Architecture", "Presentation-friendly high-level view")
    add_top_band(slide, model_name)
    labels = [
        ("Input Image", "Street scene frame"),
        ("Backbone", "Feature extraction"),
        ("Neck", "Multi-scale fusion"),
        ("Detection Head", "Object localization"),
        ("Output", "Boxes + classes"),
    ]
    positions = [Inches(0.9), Inches(3.0), Inches(5.2), Inches(7.45), Inches(9.75)]
    widths = [Inches(1.55), Inches(1.65), Inches(1.65), Inches(1.75), Inches(1.7)]
    centers = []
    for (title, subtitle), left, width in zip(labels, positions, widths):
        add_stage_box(slide, title, subtitle, left, Inches(2.3), width, Inches(1.1), color=PANEL_ALT)
        centers.append(left + width / 2)
    for idx in range(len(centers) - 1):
        add_arrow(slide, positions[idx] + widths[idx], Inches(2.85), positions[idx + 1], Inches(2.85), color=accent)
    add_panel(slide, Inches(0.9), Inches(4.1), Inches(11.4), Inches(1.45), color=PANEL_SOFT)
    add_textbox(slide, desc, Inches(1.15), Inches(4.42), Inches(10.9), Inches(0.8), size=17, color=TEXT)
    render_badges_text = {
        "YOLOv5s": [("Small model", ACCENT), ("Stable baseline", SUCCESS), ("Best final performer", ACCENT_2)],
        "YOLOv8n": [("Nano model", ACCENT), ("Modern design", SUCCESS), ("Strong precision", ACCENT_2)],
        "YOLOv11n": [("Nano model", ACCENT), ("Lightweight", SUCCESS), ("Comparison model", ACCENT_2)],
    }
    left = Inches(1.0)
    for text, badge_color in render_badges_text.get(model_name, []):
        shape = add_panel(slide, left, Inches(6.0), Inches(1.8), Inches(0.45), color=badge_color)
        shape.line.color.rgb = badge_color
        add_textbox(slide, text, left + Inches(0.05), Inches(6.08), Inches(1.7), Inches(0.25), size=11, color=BG, bold=True, align=PP_ALIGN.CENTER)
        left += Inches(2.0)
    add_footer(slide)


def build_presentation() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    baseline_df = pd.read_csv(RESULTS_ROOT / "comparison" / "baseline_comparison.csv")
    adapted_df = pd.read_csv(RESULTS_ROOT / "comparison" / "yolov5s_baseline_vs_adapted.csv")

    clear_sample = RESULTS_ROOT / "eda" / "val_clear_samples.png"
    foggy_sample = RESULTS_ROOT / "eda" / "val_foggy_samples.png"
    compare_plot = RESULTS_ROOT / "comparison" / "baseline_comparison_combined.png"
    adapt_plot = RESULTS_ROOT / "comparison" / "yolov5s_baseline_vs_adapted.png"
    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Final Project Presentation")
    add_panel(slide, Inches(0.55), Inches(0.7), Inches(12.2), Inches(5.55), color=PANEL)
    add_textbox(
        slide,
        "Object Detection in Clear and Foggy Urban Scenes Using YOLO-Based Baseline and Adaptation Models",
        Inches(0.95),
        Inches(1.05),
        Inches(10.8),
        Inches(1.4),
        size=26,
        bold=True,
    )
    add_textbox(slide, "Final Year Project Presentation", Inches(0.95), Inches(2.15), Inches(4.2), Inches(0.4), size=16, color=ACCENT, bold=True)
    add_textbox(
        slide,
        "Abizer Jesawada  |  PRN: 25070149003\nBVS Supriya  |  PRN: 25070149004",
        Inches(0.95),
        Inches(3.0),
        Inches(4.8),
        Inches(1.0),
        size=19,
    )
    add_kpi_row(
        slide,
        [
            ("Best Baseline", "YOLOv5s"),
            ("mAP50", "0.510"),
            ("mAP50-95", "0.311"),
            ("Target Domain", "Foggy Cityscapes"),
        ],
        4.65,
    )
    add_footer(slide)

    # Slide 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Problem Context")
    add_title(slide, "Problem Statement", "Why this project matters")
    add_panel(slide, Inches(0.65), Inches(1.35), Inches(12.0), Inches(4.9))
    add_bullets(
        slide,
        [
            "Object detection in urban driving scenes is important for perception systems.",
            "Models trained on clear-weather images often perform worse on foggy images.",
            "Fog reduces visibility, weakens edges, and hides distant objects.",
            "This creates a domain shift between source and target data.",
        ],
        Inches(0.95),
        Inches(1.85),
        Inches(11.2),
        Inches(3.7),
        font_size=20,
    )
    add_kpi_row(
        slide,
        [
            ("Source Domain", "Clear Images"),
            ("Target Domain", "Foggy Images"),
            ("Task", "7-Class Detection"),
            ("Core Challenge", "Domain Shift"),
        ],
        5.85,
    )
    add_footer(slide)

    # Slide 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Project Targets")
    add_title(slide, "Objectives", "Project targets")
    add_panel(slide, Inches(0.7), Inches(1.4), Inches(12.0), Inches(4.9))
    add_bullets(
        slide,
        [
            "Prepare Cityscapes and Foggy Cityscapes for YOLO-based detection.",
            "Train and compare multiple lightweight YOLO baseline models.",
            "Study model behavior on foggy target-domain images.",
            "Generate pseudo-labels and run one pilot self-training adaptation experiment.",
            "Build a professional Streamlit GUI for visual comparison and result review.",
        ],
        Inches(1.0),
        Inches(1.85),
        Inches(11.2),
        Inches(3.8),
        font_size=20,
    )
    add_footer(slide)

    # Slide 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Dataset Samples")
    add_title(slide, "Dataset Overview", "Source and target domains")
    add_metric_box(slide, "Source Dataset", "Cityscapes", Inches(0.8), Inches(1.3), Inches(2.6), Inches(0.95))
    add_metric_box(slide, "Target Dataset", "Foggy Cityscapes", Inches(3.6), Inches(1.3), Inches(2.8), Inches(0.95), accent=SUCCESS)
    add_metric_box(slide, "Detection Classes", "7 classes", Inches(6.7), Inches(1.3), Inches(2.2), Inches(0.95), accent=ACCENT_2)
    add_metric_box(slide, "Challenge", "Clear to Foggy Shift", Inches(9.2), Inches(1.3), Inches(3.0), Inches(0.95), accent=ACCENT)
    add_image(slide, clear_sample, Inches(0.8), Inches(2.6), width=Inches(5.6), height=Inches(2.9))
    add_image(slide, foggy_sample, Inches(6.9), Inches(2.6), width=Inches(5.4), height=Inches(2.9))
    add_textbox(slide, "Selected classes: person, rider, car, truck, bus, motorcycle, bicycle", Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.4), size=16, color=MUTED)
    add_footer(slide)

    # Slide 5 workflow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Pipeline")
    add_title(slide, "Project Workflow", "End-to-end pipeline")
    stages = [
        ("1", "Data Audit", "Check folders, labels, and splits"),
        ("2", "Preprocessing", "Filter classes and validate labels"),
        ("3", "YOLO Conversion", "Build YOLO txt labels and YAML"),
        ("4", "Baseline Training", "Train YOLOv5s, YOLOv8n, YOLOv11n"),
        ("5", "Evaluation", "Compare precision, recall, mAP"),
        ("6", "Foggy Inference", "Run models on foggy target images"),
        ("7", "Pseudo-Labels", "Save model predictions as labels"),
        ("8", "Self-Training", "Retrain using clear + pseudo-labeled foggy data"),
        ("9", "Final Output", "Compare results and deploy Streamlit GUI"),
    ]
    x = Inches(0.65)
    y = Inches(1.75)
    for idx, (num, title, subtitle) in enumerate(stages):
        width = Inches(1.22)
        box = add_panel(slide, x, y, width, Inches(2.2), color=PANEL_ALT if idx < 5 else PANEL_SOFT)
        add_textbox(slide, num, x + Inches(0.44), y + Inches(0.1), Inches(0.35), Inches(0.35), size=18, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, title, x + Inches(0.08), y + Inches(0.55), Inches(1.06), Inches(0.45), size=13, bold=True, align=PP_ALIGN.CENTER)
        add_textbox(slide, subtitle, x + Inches(0.08), y + Inches(1.1), Inches(1.06), Inches(0.8), size=9, color=MUTED, align=PP_ALIGN.CENTER)
        if idx < len(stages) - 1:
            add_arrow(slide, x + width, y + Inches(1.1), x + width + Inches(0.18), y + Inches(1.1), color=ACCENT)
        x += Inches(1.4)
    add_footer(slide)

    # Slide 6 preprocessing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Data Preparation")
    add_title(slide, "Preprocessing and Annotation Conversion", "Preparing reliable YOLO training data")
    add_panel(slide, Inches(0.7), Inches(1.4), Inches(6.0), Inches(5.1))
    add_bullets(
        slide,
        [
            "Converted original annotations into YOLO text label format.",
            "Filtered the project to 7 road-object detection classes.",
            "Ensured image and label filenames matched exactly.",
            "Created dataset YAML files for Ultralytics training.",
            "Validated labels using script checks and visual overlays.",
        ],
        Inches(1.0),
        Inches(1.85),
        Inches(5.3),
        Inches(4.0),
        font_size=18,
    )
    add_panel(slide, Inches(7.0), Inches(1.4), Inches(5.6), Inches(5.1), color=PANEL_ALT)
    add_textbox(slide, "Key preprocessing checks", Inches(7.3), Inches(1.8), Inches(2.5), Inches(0.3), size=17, color=ACCENT, bold=True)
    add_bullets(
        slide,
        [
            "Correct image-label filename matching",
            "YOLO normalized box coordinates",
            "Class filtering for 7 selected categories",
            "Visual verification of generated labels",
            "Ultralytics-ready dataset YAML",
        ],
        Inches(7.25),
        Inches(2.25),
        Inches(5.0),
        Inches(2.8),
        font_size=17,
    )
    add_metric_box(slide, "Outcome", "Reliable training dataset", Inches(7.35), Inches(5.45), Inches(4.7), Inches(0.85), accent=SUCCESS)
    add_footer(slide)

    # Slide 7 baseline models
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Baseline Strategy")
    add_title(slide, "Baseline Approach", "Source-only training on clear Cityscapes")
    add_metric_box(slide, "Model 1", "YOLOv5s", Inches(0.85), Inches(1.4), Inches(2.4), Inches(1.0))
    add_metric_box(slide, "Model 2", "YOLOv8n", Inches(3.55), Inches(1.4), Inches(2.4), Inches(1.0), accent=SUCCESS)
    add_metric_box(slide, "Model 3", "YOLOv11n", Inches(6.25), Inches(1.4), Inches(2.4), Inches(1.0), accent=ACCENT_2)
    add_metric_box(slide, "Training Mode", "Source Only", Inches(8.95), Inches(1.4), Inches(3.0), Inches(1.0), accent=ACCENT)
    add_bullets(
        slide,
        [
            "Each baseline model was trained only on clear-weather labeled images.",
            "These models gave the starting reference results for the project.",
            "No foggy training data was used at this stage.",
            "The best baseline model was selected for the later adaptation attempt.",
        ],
        Inches(0.9),
        Inches(3.0),
        Inches(5.4),
        Inches(2.7),
        font_size=18,
    )
    add_image(slide, compare_plot, Inches(6.7), Inches(2.7), width=Inches(5.3), height=Inches(3.0))
    add_footer(slide)

    # Slide 8 baseline table
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Results")
    add_title(slide, "Baseline Model Comparison", "Quantitative validation results")
    table_df = baseline_df.copy()
    add_table(slide, table_df, Inches(0.75), Inches(1.55), Inches(7.2), Inches(2.0))
    add_image(slide, compare_plot, Inches(8.3), Inches(1.6), width=Inches(4.4), height=Inches(4.8))
    add_textbox(slide, "YOLOv5s achieved the best overall balance of recall and mAP, making it the strongest baseline for the next stage.", Inches(0.9), Inches(4.3), Inches(6.9), Inches(1.0), size=20)
    add_footer(slide)

    # Slide 9 pseudo labels
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Pseudo-Labels")
    add_title(slide, "Pseudo-Label Generation", "Creating temporary labels for foggy images")
    add_workflow_box(slide, "Foggy Image", Inches(0.8), Inches(2.5), width=Inches(1.8), height=Inches(0.9), color=PANEL_ALT)
    add_workflow_box(slide, "Baseline Model\nPrediction", Inches(3.1), Inches(2.5), width=Inches(2.1), height=Inches(0.9), color=PANEL_ALT)
    add_workflow_box(slide, "Saved Pseudo-\nLabel .txt", Inches(5.9), Inches(2.5), width=Inches(2.0), height=Inches(0.9), color=PANEL_ALT)
    add_workflow_box(slide, "Used in\nAdaptation", Inches(8.6), Inches(2.5), width=Inches(1.8), height=Inches(0.9), color=PANEL_ALT)
    add_arrow(slide, Inches(2.6), Inches(2.95), Inches(3.1), Inches(2.95))
    add_arrow(slide, Inches(5.2), Inches(2.95), Inches(5.9), Inches(2.95))
    add_arrow(slide, Inches(7.9), Inches(2.95), Inches(8.6), Inches(2.95))
    add_bullets(
        slide,
        [
            "All three baseline models were used to generate pseudo-labels on foggy Frankfurt images.",
            "YOLOv5s produced the strongest pseudo-label quality in qualitative comparison.",
            "These labels were saved in YOLO format and later used in the pilot adaptation experiment.",
        ],
        Inches(0.9),
        Inches(4.2),
        Inches(6.2),
        Inches(2.0),
        font_size=18,
    )
    add_metric_box(slide, "Selected Pseudo-Label Source", "YOLOv5s", Inches(8.15), Inches(4.35), Inches(3.9), Inches(0.95), accent=SUCCESS)
    add_textbox(slide, "Reason: strongest baseline performance and cleaner foggy-scene predictions.", Inches(8.2), Inches(5.45), Inches(4.0), Inches(0.7), size=16, color=TEXT)
    add_footer(slide)

    # Slide 10 adaptation
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Adaptive Training")
    add_title(slide, "Adaptive / Self-Training Approach", "Where self-training happens")
    add_bullets(
        slide,
        [
            "Baseline YOLOv5s first learned from clear labeled Cityscapes images.",
            "The same baseline model then generated pseudo-labels on foggy Frankfurt images.",
            "Those pseudo-labels were combined with clear labeled data for another training run.",
            "This retraining step is the self-training part of the project.",
        ],
        Inches(0.9),
        Inches(1.7),
        Inches(5.6),
        Inches(3.4),
        font_size=19,
    )
    add_workflow_box(slide, "Clear Labeled Data", Inches(6.7), Inches(1.8), width=Inches(2.0), height=Inches(0.9), color=PANEL_ALT)
    add_workflow_box(slide, "Foggy Pseudo-\nLabeled Data", Inches(9.5), Inches(1.8), width=Inches(2.0), height=Inches(0.9), color=PANEL_ALT)
    add_workflow_box(slide, "Adapted\nYOLOv5s", Inches(8.1), Inches(3.3), width=Inches(2.1), height=Inches(1.0), color=PANEL)
    add_arrow(slide, Inches(7.7), Inches(2.7), Inches(8.7), Inches(3.3), color=SUCCESS)
    add_arrow(slide, Inches(10.5), Inches(2.7), Inches(9.5), Inches(3.3), color=SUCCESS)
    add_metric_box(slide, "Self-Training Meaning", "Model learns from its own pseudo-labels", Inches(6.85), Inches(4.75), Inches(5.1), Inches(0.95), accent=ACCENT_2)
    add_footer(slide)

    # Slide 11 adaptation results
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Pilot Experiment")
    add_title(slide, "Adaptation Experiment Results", "Baseline vs adapted YOLOv5s")
    add_table(slide, adapted_df, Inches(0.8), Inches(1.7), Inches(6.0), Inches(1.4))
    add_image(slide, adapt_plot, Inches(7.15), Inches(1.6), width=Inches(5.0), height=Inches(3.8))
    add_bullets(
        slide,
        [
            "Adaptation increased precision slightly but reduced recall.",
            "Overall mAP50 and mAP50-95 became lower than the baseline.",
            "So the pilot adaptation did not outperform baseline YOLOv5s in this setup.",
        ],
        Inches(0.9),
        Inches(4.1),
        Inches(6.0),
        Inches(1.8),
        font_size=18,
    )
    add_footer(slide)

    # Slide 12 system architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Architecture")
    add_title(slide, "Main System Architecture", "Overall project design")
    top_nodes = [
        ("Input Datasets", "Cityscapes + Foggy Cityscapes", Inches(0.75)),
        ("Data Preparation", "Audit, conversion, validation", Inches(3.15)),
        ("Baseline Models", "YOLOv5s / YOLOv8n / YOLOv11n", Inches(5.65)),
        ("Target Inference", "Foggy image prediction", Inches(8.25)),
        ("Pseudo-Labels", "Temporary target labels", Inches(10.55)),
    ]
    centers = []
    for title, sub, left in top_nodes:
        add_stage_box(slide, title, sub, left, Inches(1.95), Inches(1.95), Inches(1.2), color=PANEL_ALT)
        centers.append(left + Inches(0.975))
    for idx in range(len(centers) - 1):
        add_arrow(slide, top_nodes[idx][2] + Inches(1.95), Inches(2.55), top_nodes[idx + 1][2], Inches(2.55), color=ACCENT)
    add_stage_box(slide, "Self-Training Adaptation", "Retrain with clear + pseudo-labeled foggy data", Inches(3.0), Inches(4.45), Inches(3.1), Inches(1.25), color=PANEL)
    add_stage_box(slide, "Final Comparison", "Baseline vs adapted results", Inches(6.1), Inches(4.45), Inches(2.5), Inches(1.25), color=PANEL)
    add_stage_box(slide, "Streamlit GUI", "Interactive result review", Inches(9.0), Inches(4.45), Inches(2.5), Inches(1.25), color=PANEL)
    add_arrow(slide, Inches(11.52), Inches(3.15), Inches(4.6), Inches(4.45), color=ACCENT_2)
    add_arrow(slide, Inches(6.08), Inches(5.08), Inches(6.98), Inches(5.08), color=SUCCESS)
    add_arrow(slide, Inches(8.6), Inches(5.08), Inches(9.0), Inches(5.08), color=SUCCESS)
    add_footer(slide)

    # Slide 13-15 model architecture
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_architecture(slide, "YOLOv5s", "YOLOv5s is a small, stable, and practical detector. In this project it delivered the best overall baseline performance, so it was also selected for pseudo-label generation and the pilot self-training adaptation experiment.", ACCENT)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_architecture(slide, "YOLOv8n", "YOLOv8n is a nano model with a lightweight modern design. In this project it provided strong precision and served as a useful compact baseline and pseudo-label generator for foggy-scene comparison.", SUCCESS)

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_architecture(slide, "YOLOv11n", "YOLOv11n is a lightweight newer-generation nano detector. In this project it was used as a comparison baseline and for pseudo-label generation, but it produced the lowest overall final mAP among the three completed models.", ACCENT_2)

    # Slide 16 GUI
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "GUI")
    add_title(slide, "Streamlit GUI", "Interactive model comparison dashboard")
    features = [
        ("Upload Input Image", "Test one image directly from the interface."),
        ("Choose Models", "Select a baseline and an adapted weight file."),
        ("Side-by-Side Prediction", "Compare outputs on the same page."),
        ("Stored Metrics", "Review evaluation values and detection summary."),
        ("Project Dashboard", "Open experiment plots and overall comparison view."),
    ]
    x_positions = [Inches(0.9), Inches(4.4), Inches(7.9)]
    y_positions = [Inches(1.8), Inches(4.05)]
    idx = 0
    for y in y_positions:
        for x in x_positions:
            if idx >= len(features):
                break
            title, sub = features[idx]
            add_stage_box(slide, title, sub, x, y, Inches(3.0), Inches(1.55), color=PANEL_ALT if idx % 2 == 0 else PANEL)
            idx += 1
    add_footer(slide)

    # Slide 17 conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_top_band(slide, "Final Summary")
    add_title(slide, "Conclusion and Future Work", "Final takeaways")
    add_panel(slide, Inches(0.7), Inches(1.5), Inches(5.9), Inches(4.8))
    add_textbox(slide, "Conclusion", Inches(0.95), Inches(1.8), Inches(2.0), Inches(0.3), size=18, color=ACCENT, bold=True)
    add_bullets(
        slide,
        [
            "The project delivered a full clear-to-foggy detection workflow.",
            "YOLOv5s was the best baseline model.",
            "Pseudo-label generation worked successfully on foggy images.",
            "Pilot self-training adaptation did not improve over the baseline.",
            "The GUI made the results easier to demonstrate and compare.",
        ],
        Inches(0.95),
        Inches(2.2),
        Inches(5.2),
        Inches(3.4),
        font_size=17,
    )
    add_panel(slide, Inches(6.85), Inches(1.5), Inches(5.8), Inches(4.8), color=PANEL_ALT)
    add_textbox(slide, "Future Work", Inches(7.1), Inches(1.8), Inches(2.0), Inches(0.3), size=18, color=ACCENT_2, bold=True)
    add_bullets(
        slide,
        [
            "Try multi-city foggy adaptation instead of only Frankfurt.",
            "Tune pseudo-label confidence thresholds.",
            "Adapt YOLOv8n and YOLOv11n for broader comparison.",
            "Use stronger pseudo-label filtering strategies.",
            "Expand the GUI with downloads and more analytics.",
        ],
        Inches(7.1),
        Inches(2.2),
        Inches(5.0),
        Inches(3.4),
        font_size=17,
    )
    add_footer(slide)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Saved presentation to: {OUTPUT_PATH}")


if __name__ == "__main__":
    build_presentation()
